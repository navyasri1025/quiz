import os
import json
import re
import uuid
import logging
import traceback
from datetime import datetime
from flask import Flask, jsonify, request
from flask_cors import CORS
from dotenv import load_dotenv
from pptx import Presentation
import requests

# Configure detailed logging
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s [%(levelname)s] %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)
logger = logging.getLogger(__name__)

load_dotenv()

app = Flask(__name__)
CORS(app, resources={r"/api/*": {"origins": "*"}})

app.config["UPLOAD_FOLDER"] = os.path.join(os.path.dirname(__file__), "uploads")
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024  # 50MB
os.makedirs(app.config["UPLOAD_FOLDER"], exist_ok=True)

# In-memory storage (replace with DB in production)
quiz_sessions = {}
quiz_history = []


def extract_text_from_pptx(filepath):
    """Extract text content from PPTX file using python-pptx."""
    logger.info(f"Extracting text from PPTX: {filepath}")
    try:
        if not os.path.exists(filepath):
            raise FileNotFoundError(f"File not found: {filepath}")
        if os.path.getsize(filepath) == 0:
            raise ValueError("Uploaded file is empty")

        prs = Presentation(filepath)
        slides_text = []
        total_words = 0
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                slide_text.append(cell.text.strip())
            
            combined = "\n".join(slide_text)
            words = len(combined.split())
            total_words += words
            slides_text.append({
                "slide_number": i,
                "text": combined,
                "word_count": words
            })
        
        logger.info(f"Extracted {total_words} words from {len(slides_text)} slides")
        return {
            "slide_count": len(slides_text),
            "total_words": total_words,
            "slides": slides_text
        }
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {str(e)}")
        raise Exception(f"Error extracting text from PPTX: {str(e)}")


def generate_quiz_with_ai(content_text, question_count, difficulty):
    """Generate quiz questions using OpenRouter DeepSeek API."""
    api_key = os.getenv("OPENROUTER_API_KEY")
    if not api_key:
        logger.error("OPENROUTER_API_KEY not configured in .env")
        raise Exception("OPENROUTER_API_KEY not configured. Please add your OpenRouter API key to the .env file.")
    
    logger.info(f"Generating {question_count} {difficulty} questions with AI")
    
    difficulty_prompts = {
        "easy": "Generate EASY questions covering basic recall and simple understanding.",
        "medium": "Generate MEDIUM questions requiring comprehension and application.",
        "hard": "Generate HARD questions requiring analysis, evaluation, and synthesis."
    }
    
    prompt = f"""You are a professional quiz generator. Based on the following presentation content, generate {question_count} multiple-choice questions at {difficulty.upper()} difficulty level.

{difficulty_prompts.get(difficulty, difficulty_prompts["medium"])}

Presentation Content:
{content_text}

Return ONLY a valid JSON array (no markdown, no code fences). Each question object must have this exact structure:
{{
  "question": "The question text",
  "options": ["Option A", "Option B", "Option C", "Option D"],
  "correctAnswer": 0,
  "explanation": "Brief explanation of why this answer is correct"
}}

The correctAnswer field must be the 0-based index of the correct option (0, 1, 2, or 3).
Generate exactly {question_count} questions. Make sure options are plausible and the correct answer is accurate based on the content provided."""

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
        "HTTP-Referer": os.getenv("SITE_URL", "http://localhost:3000"),
        "X-Title": "AI Quiz Generator"
    }
    
    payload = {
        "model": "deepseek/deepseek-chat-v3-0324:free",
        "messages": [
            {"role": "system", "content": "You are a professional quiz generator that outputs only valid JSON."},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.7,
        "max_tokens": 8192
    }
    
    try:
        logger.info("Sending request to OpenRouter API...")
        response = requests.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers=headers,
            json=payload,
            timeout=120
        )
        
        logger.info(f"OpenRouter response status: {response.status_code}")
        
        if response.status_code == 401:
            raise Exception("Invalid API key. Please check your OPENROUTER_API_KEY in .env")
        if response.status_code == 429:
            raise Exception("API rate limit exceeded. Please try again later.")
        
        response.raise_for_status()
        data = response.json()
        
        if "choices" not in data or len(data["choices"]) == 0:
            logger.error(f"Unexpected API response structure: {json.dumps(data, indent=2)[:500]}")
            raise Exception("AI API returned an unexpected response format")
        
        content = data["choices"][0]["message"]["content"]
        logger.debug(f"Raw AI response: {content[:200]}...")
        
        # Clean the response - remove markdown code fences if present
        content = re.sub(r'```json\s*', '', content)
        content = re.sub(r'```\s*', '', content)
        content = content.strip()
        
        if not content:
            raise Exception("AI returned empty content")
        
        # Try to extract JSON array if wrapped in extra text
        json_match = re.search(r'\[.*\]', content, re.DOTALL)
        if json_match:
            content = json_match.group()
        
        questions = json.loads(content)
        
        if not isinstance(questions, list):
            raise ValueError("AI response is not a JSON array")
        
        if len(questions) == 0:
            raise ValueError("AI returned an empty question list")
        
        # Validate structure
        for i, q in enumerate(questions):
            if "question" not in q or "options" not in q or "correctAnswer" not in q or "explanation" not in q:
                logger.warning(f"Question {i} missing fields: {json.dumps(q)[:200]}")
                raise ValueError(f"Invalid question structure from AI at index {i}")
            if len(q["options"]) != 4:
                raise ValueError(f"Question {i} must have exactly 4 options (got {len(q['options'])})")
            if not isinstance(q["correctAnswer"], int) or q["correctAnswer"] < 0 or q["correctAnswer"] > 3:
                raise ValueError(f"Question {i} correctAnswer must be 0-3")
        
        logger.info(f"Successfully generated {len(questions)} questions")
        return questions
    
    except json.JSONDecodeError as e:
        logger.error(f"Failed to parse AI response as JSON: {content[:500]}")
        raise Exception(f"Failed to parse AI response as JSON. Raw response: {content[:200]}...")
    except requests.exceptions.Timeout:
        logger.error("AI API request timed out after 120 seconds")
        raise Exception("AI API request timed out. Please try again.")
    except requests.exceptions.RequestException as e:
        logger.error(f"API request failed: {str(e)}")
        raise Exception(f"API request failed: {str(e)}")
    except Exception as e:
        logger.error(f"Quiz generation failed: {str(e)}\n{traceback.format_exc()}")
        raise Exception(f"Quiz generation failed: {str(e)}")


@app.route("/api/health", methods=["GET"])
def health():
    """Health check endpoint."""
    logger.info("Health check requested")
    return jsonify({
        "status": "ok",
        "message": "Quiz Generator API is running",
        "timestamp": str(datetime.now())
    })


@app.route("/api/upload", methods=["POST"])
def upload_pptx():
    """Upload and process a PPT/PPTX file."""
    logger.info("Upload endpoint called")
    
    # Check if file was provided
    if "file" not in request.files:
        logger.error("No file part in request")
        return jsonify({"error": "No file provided. Please select a file to upload."}), 400
    
    file = request.files["file"]
    
    if file.filename == "":
        logger.error("Empty filename")
        return jsonify({"error": "No file selected. Please choose a file."}), 400
    
    logger.info(f"Received file: {file.filename}")
    
    if not file.filename.lower().endswith((".ppt", ".pptx")):
        logger.error(f"Invalid file format: {file.filename}")
        return jsonify({"error": "Invalid file format. Only PPT and PPTX files are allowed."}), 400
    
    # Generate unique filename
    ext = os.path.splitext(file.filename)[1]
    unique_filename = f"{uuid.uuid4()}{ext}"
    filepath = os.path.join(app.config["UPLOAD_FOLDER"], unique_filename)
    
    try:
        file.save(filepath)
        logger.info(f"File saved to: {filepath}")
        
        extracted = extract_text_from_pptx(filepath)
        
        session_id = str(uuid.uuid4())
        full_text = "\n".join([s["text"] for s in extracted["slides"]])
        
        if not full_text.strip():
            logger.warning("No text content extracted from presentation")
            # Clean up empty file
            if os.path.exists(filepath):
                os.remove(filepath)
            return jsonify({
                "error": "No text content found in the presentation. The PPTX file may contain only images."
            }), 400
        
        quiz_sessions[session_id] = {
            "filepath": filepath,
            "filename": file.filename,
            "extracted_content": extracted,
            "full_text": full_text
        }
        
        logger.info(f"Upload successful. Session ID: {session_id}, Slides: {extracted['slide_count']}")
        
        return jsonify({
            "session_id": session_id,
            "filename": file.filename,
            "slide_count": extracted["slide_count"],
            "total_words": extracted["total_words"],
            "slides": extracted["slides"]
        }), 200
    
    except Exception as e:
        logger.error(f"Upload failed: {str(e)}\n{traceback.format_exc()}")
        # Clean up on error
        if os.path.exists(filepath):
            os.remove(filepath)
        return jsonify({"error": str(e)}), 500


@app.route("/api/generate-quiz", methods=["POST"])
def generate_quiz():
    """Generate quiz questions from uploaded content."""
    logger.info("Generate quiz endpoint called")
    
    if not request.is_json:
        logger.error("Request Content-Type is not JSON")
        return jsonify({"error": "Request must be JSON. Set Content-Type: application/json."}), 400
    
    data = request.get_json()
    if not data:
        logger.error("Empty JSON body")
        return jsonify({"error": "Empty request body. Provide session_id, question_count, and difficulty."}), 400
    
    session_id = data.get("session_id")
    question_count = int(data.get("question_count", 10))
    difficulty = data.get("difficulty", "medium")
    
    logger.info(f"Session ID: {session_id}, Questions: {question_count}, Difficulty: {difficulty}")
    
    if not session_id:
        return jsonify({"error": "Missing session_id. Please upload a file first."}), 400
    
    if session_id not in quiz_sessions:
        logger.error(f"Invalid session ID: {session_id}")
        return jsonify({"error": "Invalid or expired session. Please upload a file first."}), 400
    
    if question_count < 5 or question_count > 30:
        return jsonify({"error": "Question count must be between 5 and 30."}), 400
    
    if difficulty not in ["easy", "medium", "hard"]:
        return jsonify({"error": "Difficulty must be easy, medium, or hard."}), 400
    
    session = quiz_sessions[session_id]
    full_text = session["full_text"]
    
    if not full_text.strip():
        return jsonify({"error": "No text content found in the presentation."}), 400
    
    try:
        questions = generate_quiz_with_ai(full_text, question_count, difficulty)
        
        # Store quiz in session
        session["questions"] = questions
        session["difficulty"] = difficulty
        session["question_count"] = question_count
        
        logger.info(f"Quiz generated successfully: {len(questions)} questions")
        
        return jsonify({
            "session_id": session_id,
            "questions": questions,
            "total_questions": len(questions),
            "difficulty": difficulty
        }), 200
    
    except Exception as e:
        logger.error(f"Quiz generation failed: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/submit-quiz", methods=["POST"])
def submit_quiz():
    """Submit quiz answers and get results."""
    logger.info("Submit quiz endpoint called")
    
    if not request.is_json:
        logger.error("Request Content-Type is not JSON")
        return jsonify({"error": "Request must be JSON."}), 400
    
    data = request.get_json()
    session_id = data.get("session_id") if data else None
    answers = data.get("answers", {}) if data else {}
    
    if not session_id:
        return jsonify({"error": "Missing session_id."}), 400
    
    if session_id not in quiz_sessions:
        logger.error(f"Invalid session ID for submit: {session_id}")
        return jsonify({"error": "Invalid or expired session."}), 400
    
    session = quiz_sessions[session_id]
    questions = session.get("questions", [])
    
    if not questions:
        return jsonify({"error": "No quiz found. Please generate a quiz first."}), 400
    
    results = []
    correct_count = 0
    
    for i, q in enumerate(questions):
        user_answer = answers.get(str(i))
        is_correct = user_answer == q["correctAnswer"]
        if is_correct:
            correct_count += 1
        
        results.append({
            "question_number": i + 1,
            "question": q["question"],
            "options": q["options"],
            "correct_answer": q["correctAnswer"],
            "correct_text": q["options"][q["correctAnswer"]],
            "user_answer": user_answer,
            "user_answer_text": q["options"][user_answer] if user_answer is not None else None,
            "is_correct": is_correct,
            "explanation": q["explanation"]
        })
    
    total_questions = len(questions)
    score = correct_count
    percentage = round((correct_count / total_questions) * 100, 2)
    
    result = {
        "session_id": session_id,
        "score": score,
        "total_questions": total_questions,
        "percentage": percentage,
        "passed": percentage >= 60,
        "correct_count": correct_count,
        "wrong_count": total_questions - correct_count,
        "results": results,
        "difficulty": session.get("difficulty", "medium"),
        "filename": session.get("filename", "Unknown")
    }
    
    # Save to history
    history_entry = {
        "id": str(uuid.uuid4()),
        "filename": session.get("filename", "Unknown"),
        "difficulty": session.get("difficulty", "medium"),
        "score": score,
        "total_questions": total_questions,
        "percentage": percentage,
        "passed": percentage >= 60,
        "timestamp": str(datetime.now())
    }
    quiz_history.append(history_entry)
    
    logger.info(f"Quiz submitted: Score {score}/{total_questions} ({percentage}%)")
    
    return jsonify(result), 200


@app.route("/api/history", methods=["GET"])
def get_history():
    """Get quiz history."""
    logger.info(f"History endpoint called. Total entries: {len(quiz_history)}")
    return jsonify({"history": quiz_history[-50:]}), 200


@app.route("/api/session/<session_id>", methods=["GET"])
def get_session(session_id):
    """Get session info."""
    logger.info(f"Session info requested for: {session_id}")
    if session_id not in quiz_sessions:
        return jsonify({"error": "Session not found"}), 404
    
    session = quiz_sessions[session_id]
    return jsonify({
        "filename": session.get("filename"),
        "slide_count": session.get("extracted_content", {}).get("slide_count"),
        "total_words": session.get("extracted_content", {}).get("total_words"),
        "has_questions": "questions" in session
    }), 200


# Global error handler for 404s
@app.errorhandler(404)
def not_found(e):
    logger.warning(f"404 error: {str(e)}")
    return jsonify({"error": "Endpoint not found. Available endpoints: /api/health, /api/upload, /api/generate-quiz, /api/submit-quiz, /api/history, /api/session/<id>"}), 404


# Global error handler for 413 (file too large)
@app.errorhandler(413)
def too_large(e):
    logger.warning("File too large error")
    return jsonify({"error": "File too large. Maximum size is 50MB."}), 413


# Global error handler for 500
@app.errorhandler(500)
def internal_error(e):
    logger.error(f"Internal server error: {str(e)}")
    return jsonify({"error": "Internal server error. Please try again."}), 500


if __name__ == "__main__":
    port = int(os.getenv("PORT", 5000))
    debug = os.getenv("FLASK_DEBUG", "false").lower() == "true"
    
    api_key_status = "CONFIGURED" if os.getenv("OPENROUTER_API_KEY") else "NOT SET"
    logger.info(f"Starting Quiz Generator API on port {port} (debug={debug})")
    logger.info(f"OpenRouter API Key: {api_key_status}")
    logger.info(f"Upload folder: {app.config['UPLOAD_FOLDER']}")
    
    app.run(host="0.0.0.0", port=port, debug=debug)