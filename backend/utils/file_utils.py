"""
file_utils.py — Multi-format file extraction for QuizCraft AI.

Supported formats
-----------------
  .pptx / .ppt  — python-pptx
  .pdf          — pdfplumber
  .docx         — python-docx
  .txt          — stdlib (utf-8 with latin-1 fallback)

Every extractor returns the same shape:
    {
        "slide_count": int,      # pages / sections / 1 for flat text
        "total_words": int,
        "slides": [
            { "slide_number": int, "text": str, "word_count": int }
        ]
    }

Sanitization strips prompt-injection lines and hard-caps content length.
"""
import os
import re
import logging
import time
from typing import Union

from pptx import Presentation
import pdfplumber
import docx as python_docx

from config import (
    ALLOWED_EXTENSIONS,
    FILE_TYPE_LABELS,
    MAX_CONTENT_CHARS,
    INJECTION_PATTERNS,
)

logger = logging.getLogger("quizcraft")


# ---------------------------------------------------------------------------
# Validation helpers
# ---------------------------------------------------------------------------

def is_allowed_extension(filename: str) -> bool:
    """Return True when the file extension is in ALLOWED_EXTENSIONS."""
    return get_extension(filename) in ALLOWED_EXTENSIONS


def get_extension(filename: str) -> str:
    """Return the lowercase extension including the leading dot."""
    return os.path.splitext(filename.lower())[1]


def get_file_type_label(filename: str) -> str:
    """Return the human-readable file type label (e.g. 'PDF')."""
    return FILE_TYPE_LABELS.get(get_extension(filename), "Unknown")


# ---------------------------------------------------------------------------
# Per-format extractors
# ---------------------------------------------------------------------------

def _extract_pptx(filepath: str) -> dict:
    """Extract text from a .pptx file using python-pptx."""
    try:
        prs = Presentation(filepath)
    except Exception as exc:
        raise ValueError(
            f"Cannot open file as a PowerPoint presentation. "
            f"Ensure it is a valid .pptx file. Detail: {exc}"
        ) from exc

    slides_data = []
    total_words = 0

    for idx, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            parts.append(cell.text.strip())

        combined = "\n".join(parts)
        wc = len(combined.split()) if combined.strip() else 0
        total_words += wc
        slides_data.append({"slide_number": idx, "text": combined, "word_count": wc})

    return {"slide_count": len(slides_data), "total_words": total_words, "slides": slides_data}


def _extract_pdf(filepath: str) -> dict:
    """Extract text from a .pdf file using pdfplumber."""
    try:
        with pdfplumber.open(filepath) as pdf:
            pages_data = []
            total_words = 0

            for idx, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                text = text.strip()
                wc = len(text.split()) if text else 0
                total_words += wc
                pages_data.append({"slide_number": idx, "text": text, "word_count": wc})

    except Exception as exc:
        raise ValueError(
            f"Cannot read PDF file. Ensure it is a valid, non-encrypted PDF. "
            f"Detail: {exc}"
        ) from exc

    return {"slide_count": len(pages_data), "total_words": total_words, "slides": pages_data}


def _extract_docx(filepath: str) -> dict:
    """Extract text from a .docx file using python-docx."""
    try:
        doc = python_docx.Document(filepath)
    except Exception as exc:
        raise ValueError(
            f"Cannot open Word document. Ensure it is a valid .docx file. "
            f"Detail: {exc}"
        ) from exc

    # Group paragraphs into logical sections (~20 paragraphs each) so the
    # frontend "slide preview" shows manageable chunks.
    all_paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    # Also pull text from tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    all_paragraphs.append(cell.text.strip())

    # Split into pages of ~20 paragraphs so slide_count is meaningful.
    PAGE_SIZE = 20
    pages_data = []
    total_words = 0
    chunks = [all_paragraphs[i:i + PAGE_SIZE] for i in range(0, max(len(all_paragraphs), 1), PAGE_SIZE)]

    for idx, chunk in enumerate(chunks, start=1):
        combined = "\n".join(chunk)
        wc = len(combined.split()) if combined else 0
        total_words += wc
        pages_data.append({"slide_number": idx, "text": combined, "word_count": wc})

    if not pages_data:
        pages_data = [{"slide_number": 1, "text": "", "word_count": 0}]

    return {"slide_count": len(pages_data), "total_words": total_words, "slides": pages_data}


def _extract_txt(filepath: str) -> dict:
    """Extract text from a plain .txt file."""
    # Try UTF-8 first; fall back to latin-1 for legacy files.
    for encoding in ("utf-8", "latin-1"):
        try:
            with open(filepath, "r", encoding=encoding) as fh:
                content = fh.read()
            break
        except UnicodeDecodeError:
            continue
    else:
        raise ValueError("Could not decode text file. Ensure it is UTF-8 or Latin-1 encoded.")

    lines = [ln.strip() for ln in content.splitlines() if ln.strip()]

    # Split into pages of ~30 lines.
    PAGE_SIZE = 30
    pages_data = []
    total_words = 0
    chunks = [lines[i:i + PAGE_SIZE] for i in range(0, max(len(lines), 1), PAGE_SIZE)]

    for idx, chunk in enumerate(chunks, start=1):
        combined = "\n".join(chunk)
        wc = len(combined.split()) if combined else 0
        total_words += wc
        pages_data.append({"slide_number": idx, "text": combined, "word_count": wc})

    if not pages_data:
        pages_data = [{"slide_number": 1, "text": "", "word_count": 0}]

    return {"slide_count": len(pages_data), "total_words": total_words, "slides": pages_data}


# ---------------------------------------------------------------------------
# Unified public extractor
# ---------------------------------------------------------------------------

def extract_text(filepath: str, filename: str) -> dict:
    """
    Route to the correct extractor based on file extension.

    Returns the standard dict { slide_count, total_words, slides[] } plus
    the file_processing_time_seconds field for analytics.

    Raises FileNotFoundError, ValueError, or Exception on failure.
    """
    if not os.path.exists(filepath):
        raise FileNotFoundError(f"Saved file not found on disk: {filepath}")
    if os.path.getsize(filepath) == 0:
        raise ValueError("Uploaded file is empty (0 bytes).")

    ext = get_extension(filename)
    start = time.monotonic()

    logger.info("Extracting text from %r (ext=%s)", filename, ext)

    if ext in (".ppt", ".pptx"):
        result = _extract_pptx(filepath)
    elif ext == ".pdf":
        result = _extract_pdf(filepath)
    elif ext == ".docx":
        result = _extract_docx(filepath)
    elif ext == ".txt":
        result = _extract_txt(filepath)
    else:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            "Supported: .pptx, .ppt, .pdf, .docx, .txt"
        )

    elapsed = round(time.monotonic() - start, 3)
    result["file_processing_time_seconds"] = elapsed

    logger.info(
        "Extraction done: %d sections, %d words, %.2fs",
        result["slide_count"], result["total_words"], elapsed
    )
    return result


# Backwards-compatible alias used by existing tests / imports.
extract_text_from_pptx = extract_text


# ---------------------------------------------------------------------------
# Sanitization
# ---------------------------------------------------------------------------

def sanitize_content(raw_text: str) -> str:
    """
    Clean extracted text before inserting it into an AI prompt.

    1. Strip lines containing known prompt-injection phrases.
    2. Collapse multiple blank lines → single blank line.
    3. Truncate to MAX_CONTENT_CHARS.
    """
    lines = raw_text.splitlines()
    clean_lines = []

    for line in lines:
        if any(pat in line.lower() for pat in INJECTION_PATTERNS):
            logger.warning("Stripped injection line: %r", line[:80])
            continue
        clean_lines.append(line)

    cleaned = re.sub(r"\n{3,}", "\n\n", "\n".join(clean_lines))

    if len(cleaned) > MAX_CONTENT_CHARS:
        logger.info("Content truncated %d → %d chars", len(cleaned), MAX_CONTENT_CHARS)
        cleaned = cleaned[:MAX_CONTENT_CHARS]

    return cleaned.strip()
